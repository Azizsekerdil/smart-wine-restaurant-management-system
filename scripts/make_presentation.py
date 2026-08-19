"""Wine House tanıtım sunumu üreteci.

Tek bir içerik kaynağından (``presentation_content.py``) dört çıktı üretir:

    Wine_House_Tanitim.html          Kaydırmalı, kendi kendine yeten HTML
    Wine_House_Tanitim.pptx          PowerPoint (ekran · koyu tema)
    Wine_House_Tanitim.pdf           PDF (ekran · koyu tema)
    Wine_House_Tanitim_Baski.pptx    PowerPoint (baskı · açık tema)
    Wine_House_Tanitim_Baski.pdf     PDF (baskı · açık tema)

    Wine_House_Intro_EN.html / .pptx / .pdf
    Wine_House_Intro_EN_Print.pptx / .pdf

Kullanım::

    .venv\\Scripts\\python.exe scripts\\make_presentation.py

Notlar
------
* Slayt boyutu 13,333 × 7,5 inç (16:9).
* PDF'te Türkçe karakterler için sistemdeki bir TrueType yazı tipi gömülür;
  bulunamazsa **sessizce bozuk çıktı üretmek yerine** açık hata verilir.
* HTML tek dosyadır ve hiçbir dış kaynağa (CDN, yazı tipi, görsel) bağlı
  değildir — çevrimdışı açılır.
"""

from __future__ import annotations

import contextlib
import html as html_mod
import re
import sys
from dataclasses import dataclass
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT / "scripts"))

from presentation_content import DECKS  # noqa: E402

# ---------------------------------------------------------------------------
# Tema
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class Theme:
    """Bir sunum temasının renk paleti."""

    name: str
    bg: str
    surface: str
    surface_alt: str
    border: str
    text: str
    muted: str
    accent: str
    wine: str
    is_dark: bool


SCREEN = Theme(
    name="screen",
    bg="1A0E14",
    surface="241419",
    surface_alt="2E1A21",
    border="3D2530",
    text="F2EAEC",
    muted="B9A7AE",
    accent="C8A45C",
    wine="9B4B6B",
    is_dark=True,
)

PRINT = Theme(
    name="print",
    bg="FFFFFF",
    surface="F7F3F4",
    surface_alt="EFE7EA",
    border="D9CCD1",
    text="23181C",
    muted="6B5C62",
    accent="8A6320",
    wine="4A1C2E",
    is_dark=False,
)

#: PPTX metin yazı tipi. Depoda paketli DejaVu Sans kullanılır; tescilli bir
#: sistem yazı tipine bağımlılık yoktur (bkz. THIRD_PARTY_NOTICES.md).
PPTX_FONT = "DejaVu Sans"

SLIDE_W_IN = 13.3333
SLIDE_H_IN = 7.5


# ---------------------------------------------------------------------------
# Metin yardımcıları
# ---------------------------------------------------------------------------
def split_bold(text: str) -> list[tuple[str, bool]]:
    """``**kalın**`` işaretlerini (parça, kalın_mı) çiftlerine ayırır."""
    parts: list[tuple[str, bool]] = []
    for index, chunk in enumerate(re.split(r"\*\*(.+?)\*\*", text)):
        if chunk:
            parts.append((chunk, index % 2 == 1))
    return parts or [(text, False)]


def plain(text: str) -> str:
    """Kalın işaretlerini kaldırır."""
    return re.sub(r"\*\*(.+?)\*\*", r"\1", text)


def to_html(text: str) -> str:
    """Kalın işaretlerini HTML'e çevirir ve kaçış uygular."""
    out = []
    for chunk, bold in split_bold(text):
        escaped = html_mod.escape(chunk)
        out.append(f"<strong>{escaped}</strong>" if bold else escaped)
    return "".join(out)


def to_rl(text: str) -> str:
    """Kalın işaretlerini ReportLab imine çevirir."""
    out = []
    for chunk, bold in split_bold(text):
        escaped = (
            chunk.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        )
        out.append(f"<b>{escaped}</b>" if bold else escaped)
    return "".join(out)


# ===========================================================================
# PPTX
# ===========================================================================
def build_pptx(deck: dict, theme: Theme, target: Path) -> None:
    """PowerPoint sunumu üretir."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    def rgb(hex_value: str) -> RGBColor:
        return RGBColor.from_string(hex_value)

    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_W_IN)
    presentation.slide_height = Inches(SLIDE_H_IN)
    blank = presentation.slide_layouts[6]

    presentation.core_properties.title = deck["deck_title"]
    presentation.core_properties.author = "Aziz Şekerdil"
    presentation.core_properties.subject = deck["tagline"]

    def new_slide():
        slide = presentation.slides.add_slide(blank)
        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = rgb(theme.bg)
        return slide

    def rect(slide, x, y, w, h, fill=None, line=None, radius=False):
        from pptx.enum.shapes import MSO_SHAPE

        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, x, y, w, h)
        if radius:
            # Bazı şekil türlerinde ayar tutamacı yoktur; köşe yarıçapı atlanır.
            with contextlib.suppress(IndexError, KeyError):
                shape.adjustments[0] = 0.06
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(fill)
        else:
            shape.fill.background()
        if line:
            shape.line.color.rgb = rgb(line)
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
        box = slide.shapes.add_textbox(x, y, w, h)
        frame = box.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = anchor
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        return frame

    def write(paragraph, text, size, color, bold=False, space_after=0, align=PP_ALIGN.LEFT):
        paragraph.alignment = align
        paragraph.space_after = Pt(space_after)
        for chunk, is_bold in split_bold(text):
            run = paragraph.add_run()
            run.text = chunk
            run.font.size = Pt(size)
            run.font.bold = bold or is_bold
            run.font.color.rgb = rgb(theme.accent if is_bold and not bold else color)
            run.font.name = PPTX_FONT
        return paragraph

    def header(slide, title, subtitle=""):
        """Standart slayt başlığı ve altın ayraç."""
        frame = textbox(slide, Inches(0.85), Inches(0.62), Inches(11.6), Inches(1.0))
        write(frame.paragraphs[0], title, 30, theme.text, bold=True)
        if subtitle:
            paragraph = frame.add_paragraph()
            write(paragraph, subtitle, 14, theme.muted)
        rect(slide, Inches(0.85), Inches(1.72), Inches(1.5), Pt(3), fill=theme.accent)

    def footer(slide, index):
        frame = textbox(slide, Inches(0.85), Inches(6.92), Inches(8.0), Inches(0.35))
        write(frame.paragraphs[0], deck["footer"], 9, theme.muted)
        num = textbox(slide, Inches(11.4), Inches(6.92), Inches(1.05), Inches(0.35))
        write(num.paragraphs[0], str(index), 10, theme.muted, align=PP_ALIGN.RIGHT)

    total = len(deck["slides"])

    for index, spec in enumerate(deck["slides"], start=1):
        kind = spec["type"]
        slide = new_slide()

        # ---------------- KAPAK ----------------
        if kind in {"title", "closing"}:
            rect(slide, 0, 0, Inches(0.28), Inches(SLIDE_H_IN), fill=theme.accent)
            rect(
                slide,
                Inches(7.9),
                Inches(-1.2),
                Inches(7.0),
                Inches(10.0),
                fill=theme.surface if theme.is_dark else theme.surface_alt,
            )

            frame = textbox(slide, Inches(1.1), Inches(2.15), Inches(7.2), Inches(1.4))
            write(frame.paragraphs[0], spec["title"], 60, theme.accent, bold=True)

            sub = textbox(slide, Inches(1.1), Inches(3.45), Inches(7.0), Inches(0.9))
            write(sub.paragraphs[0], spec["subtitle"], 21, theme.text)

            if spec.get("note"):
                note = textbox(slide, Inches(1.1), Inches(4.35), Inches(7.0), Inches(0.5))
                write(note.paragraphs[0], spec["note"], 13, theme.muted)

            if kind == "closing" and spec.get("items"):
                item_frame = textbox(slide, Inches(8.55), Inches(2.1), Inches(4.1), Inches(3.4))
                for position, item in enumerate(spec["items"]):
                    paragraph = (
                        item_frame.paragraphs[0] if position == 0 else item_frame.add_paragraph()
                    )
                    write(paragraph, f"◆  {plain(item)}", 12.5, theme.text, space_after=11)

            rect(slide, Inches(1.1), Inches(5.15), Inches(2.2), Pt(3), fill=theme.accent)
            meta = textbox(slide, Inches(1.1), Inches(5.45), Inches(7.0), Inches(0.5))
            write(meta.paragraphs[0], spec.get("meta", ""), 12, theme.muted)

            footer(slide, index)
            continue

        # ---------------- BÖLÜM AYRACI ----------------
        if kind == "section":
            rect(slide, 0, Inches(2.75), Inches(SLIDE_W_IN), Inches(2.0), fill=theme.surface)
            rect(slide, Inches(0.85), Inches(2.75), Pt(5), Inches(2.0), fill=theme.accent)

            frame = textbox(slide, Inches(1.3), Inches(3.1), Inches(10.5), Inches(0.9))
            write(frame.paragraphs[0], spec["title"], 40, theme.text, bold=True)
            if spec.get("subtitle"):
                sub = textbox(slide, Inches(1.3), Inches(3.95), Inches(10.5), Inches(0.6))
                write(sub.paragraphs[0], spec["subtitle"], 16, theme.accent)

            counter = textbox(slide, Inches(0.85), Inches(1.6), Inches(4.0), Inches(0.4))
            write(counter.paragraphs[0], f"{index} / {total}", 11, theme.muted)
            footer(slide, index)
            continue

        header(slide, spec["title"], spec.get("subtitle", ""))

        # ---------------- MADDELER ----------------
        if kind == "bullets":
            items = spec["items"]
            top = 2.15
            gap = min(0.72, 4.35 / max(1, len(items)))
            for item in items:
                rect(
                    slide,
                    Inches(0.9),
                    Inches(top + gap * 0.28),
                    Pt(6),
                    Pt(6),
                    fill=theme.accent,
                )
                frame = textbox(
                    slide, Inches(1.28), Inches(top), Inches(11.2), Inches(gap)
                )
                write(frame.paragraphs[0], item, 14.5 if len(items) > 6 else 16, theme.text)
                top += gap

        # ---------------- SAYI KARTLARI ----------------
        elif kind == "stats":
            items = spec["items"]
            columns = 4 if len(items) > 6 else 3
            rows = (len(items) + columns - 1) // columns
            card_w = 11.6 / columns - 0.22
            card_h = 1.42 if rows > 2 else 1.65
            for position, (value, label) in enumerate(items):
                column, row = position % columns, position // columns
                x = Inches(0.85 + column * (card_w + 0.22))
                y = Inches(2.2 + row * (card_h + 0.2))
                rect(
                    slide, x, y, Inches(card_w), Inches(card_h),
                    fill=theme.surface, line=theme.border, radius=True,
                )
                rect(slide, x, y, Pt(4), Inches(card_h), fill=theme.accent)
                frame = textbox(
                    slide,
                    Inches(0.85 + column * (card_w + 0.22) + 0.24),
                    Inches(2.2 + row * (card_h + 0.2) + 0.2),
                    Inches(card_w - 0.4),
                    Inches(card_h - 0.3),
                )
                write(frame.paragraphs[0], value, 30, theme.accent, bold=True)
                write(frame.add_paragraph(), label, 11, theme.muted)
            if spec.get("note"):
                note = textbox(slide, Inches(0.85), Inches(6.32), Inches(11.6), Inches(0.55))
                write(note.paragraphs[0], spec["note"], 11, theme.muted)

        # ---------------- TABLO ----------------
        elif kind == "table":
            columns = spec["columns"]
            rows = spec["rows"]
            widths = (
                [0.30, 0.70] if len(columns) == 2
                else [0.26, 0.24, 0.50] if len(columns) == 3
                else [0.22, 0.18, 0.16, 0.44]
            )
            top = 2.2
            row_h = min(0.66, 4.0 / (len(rows) + 1))

            x = 0.85
            for position, name in enumerate(columns):
                width = 11.6 * widths[position]
                rect(slide, Inches(x), Inches(top), Inches(width - 0.06),
                     Inches(row_h), fill=theme.wine)
                frame = textbox(
                    slide, Inches(x + 0.16), Inches(top + row_h * 0.22),
                    Inches(width - 0.34), Inches(row_h),
                )
                write(frame.paragraphs[0], name, 12,
                      "FFFFFF" if theme.is_dark else "FFFFFF", bold=True)
                x += width
            top += row_h + 0.06

            for row_index, row in enumerate(rows):
                x = 0.85
                shade = theme.surface if row_index % 2 == 0 else theme.surface_alt
                for position, cell in enumerate(row):
                    width = 11.6 * widths[position]
                    rect(slide, Inches(x), Inches(top), Inches(width - 0.06),
                         Inches(row_h), fill=shade, line=theme.border)
                    frame = textbox(
                        slide, Inches(x + 0.16), Inches(top + row_h * 0.2),
                        Inches(width - 0.34), Inches(row_h),
                    )
                    write(frame.paragraphs[0], cell, 11.5,
                          theme.text if position else theme.accent,
                          bold=position == 0)
                    x += width
                top += row_h + 0.04

            if spec.get("note"):
                note = textbox(slide, Inches(0.85), Inches(min(6.35, top + 0.15)),
                               Inches(11.6), Inches(0.5))
                write(note.paragraphs[0], spec["note"], 11, theme.muted)

        # ---------------- AKIŞ ----------------
        elif kind == "flow":
            steps = spec["steps"]
            box_h = 0.82
            top = 2.35
            for position, step in enumerate(steps):
                y = Inches(top + position * (box_h + 0.16))
                rect(slide, Inches(0.85), y, Inches(0.82), Inches(box_h),
                     fill=theme.accent, radius=True)
                number = textbox(slide, Inches(0.85), Inches(top + position * (box_h + 0.16) + 0.2),
                                 Inches(0.82), Inches(0.5))
                write(number.paragraphs[0], str(position + 1), 20,
                      theme.bg if theme.is_dark else "FFFFFF",
                      bold=True, align=PP_ALIGN.CENTER)

                rect(slide, Inches(1.78), y, Inches(10.65), Inches(box_h),
                     fill=theme.surface, line=theme.border, radius=True)
                frame = textbox(
                    slide, Inches(2.05),
                    Inches(top + position * (box_h + 0.16) + 0.22),
                    Inches(10.1), Inches(box_h),
                )
                write(frame.paragraphs[0], step, 15, theme.text)

            if spec.get("note"):
                note = textbox(slide, Inches(0.85), Inches(6.3), Inches(11.6), Inches(0.62))
                write(note.paragraphs[0], spec["note"], 11, theme.muted)

        # ---------------- EKRAN GÖRÜNTÜSÜ ----------------
        elif kind == "screenshot":
            from PIL import Image

            image_path = PROJECT_ROOT / spec["image"]
            with Image.open(image_path) as image:
                aspect = image.height / image.width

            max_w, max_h = 11.3, 4.42
            width = min(max_w, max_h / aspect)
            height = width * aspect
            x = 0.85 + (11.63 - width) / 2
            y = 2.02

            rect(
                slide,
                Inches(x - 0.05),
                Inches(y - 0.05),
                Inches(width + 0.10),
                Inches(height + 0.10),
                fill=theme.surface,
                line=theme.border,
                radius=True,
            )
            slide.shapes.add_picture(
                str(image_path), Inches(x), Inches(y), width=Inches(width)
            )

            if spec.get("caption"):
                cap = textbox(
                    slide, Inches(0.85), Inches(y + height + 0.14), Inches(11.6), Inches(0.5)
                )
                write(cap.paragraphs[0], spec["caption"], 11.5, theme.muted)

        # ---------------- KARŞILAŞTIRMA ----------------
        elif kind == "compare":
            for side, (title_key, items_key, x0) in enumerate(
                [("left_title", "left", 0.85), ("right_title", "right", 7.1)]
            ):
                rect(slide, Inches(x0), Inches(2.15), Inches(5.38), Inches(4.0),
                     fill=theme.surface, line=theme.border, radius=True)
                rect(slide, Inches(x0), Inches(2.15), Inches(5.38), Pt(4),
                     fill=theme.wine if side == 0 else theme.accent)

                head = textbox(slide, Inches(x0 + 0.3), Inches(2.42),
                               Inches(4.8), Inches(0.5))
                write(head.paragraphs[0], spec[title_key], 15,
                      theme.wine if side == 0 and not theme.is_dark else theme.accent,
                      bold=True)

                top = 3.02
                for item in spec[items_key]:
                    frame = textbox(slide, Inches(x0 + 0.3), Inches(top),
                                    Inches(4.85), Inches(0.62))
                    write(frame.paragraphs[0], f"·  {item}", 12.5, theme.text)
                    top += 0.62

            if spec.get("note"):
                note = textbox(slide, Inches(0.85), Inches(6.32), Inches(11.6), Inches(0.5))
                write(note.paragraphs[0], spec["note"], 11, theme.muted)

        footer(slide, index)

    target.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(target))


# ===========================================================================
# PDF
# ===========================================================================
def build_pdf(deck: dict, theme: Theme, target: Path) -> None:
    """PDF sunumu üretir (Türkçe karakterler için TrueType gömer)."""
    from reportlab.lib.colors import HexColor
    from reportlab.pdfgen import canvas as rl_canvas

    from apps.reporting.exporters import _register_pdf_font  # noqa: PLC0415

    # Tek aile kullanılır; kalınlık boyut ve renkle ifade edilir.
    font = _register_pdf_font()

    W, H = SLIDE_W_IN * 72, SLIDE_H_IN * 72  # 960 × 540 punto
    c = rl_canvas.Canvas(str(target), pagesize=(W, H))
    c.setTitle(deck["deck_title"])
    c.setAuthor("Aziz Şekerdil")
    c.setSubject(deck["tagline"])

    col = {k: HexColor("#" + v) for k, v in {
        "bg": theme.bg, "surface": theme.surface, "surface_alt": theme.surface_alt,
        "border": theme.border, "text": theme.text, "muted": theme.muted,
        "accent": theme.accent, "wine": theme.wine, "white": "FFFFFF",
    }.items()}

    def wrap(text: str, size: float, max_w: float) -> list[str]:
        words, lines, current = plain(text).split(), [], ""
        for word in words:
            trial = f"{current} {word}".strip()
            if c.stringWidth(trial, font, size) <= max_w:
                current = trial
            else:
                if current:
                    lines.append(current)
                current = word
        if current:
            lines.append(current)
        return lines or [""]

    def draw_rich(text, x, y, size, base_color, max_w, leading=None):
        """``**kalın**`` parçaları vurgu rengiyle çizer, satır sarar."""
        leading = leading or size * 1.32
        segments = split_bold(text)
        cursor_x, cursor_y = x, y
        for chunk, is_bold in segments:
            for word in re.findall(r"\S+\s*", chunk):
                width = c.stringWidth(word, font, size)
                if cursor_x + width > x + max_w and cursor_x > x:
                    cursor_x, cursor_y = x, cursor_y - leading
                c.setFillColor(col["accent"] if is_bold else base_color)
                c.setFont(font, size)
                c.drawString(cursor_x, cursor_y, word)
                cursor_x += width
        return cursor_y

    def page_bg():
        c.setFillColor(col["bg"])
        c.rect(0, 0, W, H, stroke=0, fill=1)

    def footer(index):
        c.setFont(font, 7.5)
        c.setFillColor(col["muted"])
        c.drawString(61, 26, deck["footer"])
        c.drawRightString(W - 61, 26, str(index))

    def header(title, subtitle=""):
        c.setFont(font, 25)
        c.setFillColor(col["text"])
        c.drawString(61, H - 66, plain(title))
        if subtitle:
            c.setFont(font, 11)
            c.setFillColor(col["muted"])
            c.drawString(61, H - 86, plain(subtitle))
        c.setFillColor(col["accent"])
        c.rect(61, H - 100, 108, 2.5, stroke=0, fill=1)

    def box(x, y, w, h, fill, border=None, radius=6):
        c.setFillColor(fill)
        if border is not None:
            c.setStrokeColor(border)
            c.setLineWidth(0.7)
            c.roundRect(x, y, w, h, radius, stroke=1, fill=1)
        else:
            c.roundRect(x, y, w, h, radius, stroke=0, fill=1)

    total = len(deck["slides"])

    for index, spec in enumerate(deck["slides"], start=1):
        kind = spec["type"]
        page_bg()

        if kind in {"title", "closing"}:
            c.setFillColor(col["accent"])
            c.rect(0, 0, 20, H, stroke=0, fill=1)
            c.setFillColor(col["surface"] if theme.is_dark else col["surface_alt"])
            c.rect(569, -60, 460, 660, stroke=0, fill=1)

            c.setFont(font, 52)
            c.setFillColor(col["accent"])
            c.drawString(79, H - 210, spec["title"])

            c.setFont(font, 18)
            c.setFillColor(col["text"])
            c.drawString(79, H - 248, spec["subtitle"])

            if spec.get("note"):
                c.setFont(font, 11)
                c.setFillColor(col["muted"])
                c.drawString(79, H - 274, spec["note"])

            if kind == "closing" and spec.get("items"):
                y = H - 200
                for item in spec["items"]:
                    for line_index, line in enumerate(wrap(f"◆  {item}", 10.5, 285)):
                        c.setFont(font, 10.5)
                        c.setFillColor(col["text"])
                        c.drawString(616 if line_index == 0 else 628, y, line)
                        y -= 15
                    y -= 8

            c.setFillColor(col["accent"])
            c.rect(79, H - 330, 158, 2.5, stroke=0, fill=1)
            c.setFont(font, 10)
            c.setFillColor(col["muted"])
            c.drawString(79, H - 352, spec.get("meta", ""))
            footer(index)
            c.showPage()
            continue

        if kind == "section":
            c.setFillColor(col["surface"])
            c.rect(0, H - 342, W, 144, stroke=0, fill=1)
            c.setFillColor(col["accent"])
            c.rect(61, H - 342, 4, 144, stroke=0, fill=1)

            c.setFont(font, 34)
            c.setFillColor(col["text"])
            c.drawString(94, H - 248, spec["title"])
            if spec.get("subtitle"):
                c.setFont(font, 13)
                c.setFillColor(col["accent"])
                c.drawString(94, H - 274, spec["subtitle"])

            c.setFont(font, 9.5)
            c.setFillColor(col["muted"])
            c.drawString(61, H - 118, f"{index} / {total}")
            footer(index)
            c.showPage()
            continue

        header(spec["title"], spec.get("subtitle", ""))

        if kind == "bullets":
            items = spec["items"]
            size = 12 if len(items) > 6 else 13
            y = H - 132
            step = min(50, 300 / max(1, len(items)))
            for item in items:
                c.setFillColor(col["accent"])
                c.circle(66, y + 4, 3, stroke=0, fill=1)
                last_y = draw_rich(item, 82, y, size, col["text"], 800, leading=size * 1.35)
                y = min(y - step, last_y - step * 0.72)

        elif kind == "stats":
            items = spec["items"]
            columns = 4 if len(items) > 6 else 3
            card_w = (838 - (columns - 1) * 14) / columns
            card_h = 96 if len(items) > 6 else 112
            for position, (value, label) in enumerate(items):
                column, row = position % columns, position // columns
                x = 61 + column * (card_w + 14)
                y = H - 148 - row * (card_h + 13) - card_h
                box(x, y, card_w, card_h, col["surface"], col["border"])
                c.setFillColor(col["accent"])
                c.rect(x, y, 3.5, card_h, stroke=0, fill=1)
                c.setFont(font, 27)
                c.setFillColor(col["accent"])
                c.drawString(x + 18, y + card_h - 40, value)
                c.setFont(font, 9)
                c.setFillColor(col["muted"])
                for line_index, line in enumerate(wrap(label, 9, card_w - 32)):
                    c.drawString(x + 18, y + card_h - 58 - line_index * 12, line)
            if spec.get("note"):
                c.setFont(font, 8.5)
                c.setFillColor(col["muted"])
                for line_index, line in enumerate(wrap(spec["note"], 8.5, 838)):
                    c.drawString(61, 58 - line_index * 11, line)

        elif kind == "table":
            columns, rows = spec["columns"], spec["rows"]
            widths = (
                [0.30, 0.70] if len(columns) == 2
                else [0.26, 0.24, 0.50] if len(columns) == 3
                else [0.22, 0.18, 0.16, 0.44]
            )
            row_h = min(42, 268 / (len(rows) + 1))
            y = H - 148 - row_h

            x = 61
            for position, name in enumerate(columns):
                width = 838 * widths[position]
                c.setFillColor(col["wine"])
                c.rect(x, y, width - 4, row_h, stroke=0, fill=1)
                c.setFont(font, 10)
                c.setFillColor(col["white"])
                c.drawString(x + 12, y + row_h / 2 - 3.5, plain(name))
                x += width
            y -= row_h + 3

            for row_index, row in enumerate(rows):
                x = 61
                shade = col["surface"] if row_index % 2 == 0 else col["surface_alt"]
                for position, cell in enumerate(row):
                    width = 838 * widths[position]
                    c.setFillColor(shade)
                    c.setStrokeColor(col["border"])
                    c.setLineWidth(0.6)
                    c.rect(x, y, width - 4, row_h, stroke=1, fill=1)
                    lines = wrap(cell, 9.5, width - 26)[:2]
                    start = y + row_h / 2 + (5 if len(lines) > 1 else -3.5)
                    for line_index, line in enumerate(lines):
                        c.setFont(font, 9.5)
                        c.setFillColor(col["accent"] if position == 0 else col["text"])
                        c.drawString(x + 12, start - line_index * 11, line)
                    x += width
                y -= row_h + 2

            if spec.get("note"):
                c.setFont(font, 8.5)
                c.setFillColor(col["muted"])
                c.drawString(61, max(56, y - 6), plain(spec["note"]))

        elif kind == "flow":
            steps = spec["steps"]
            box_h = 46
            y = H - 148 - box_h
            for position, step in enumerate(steps):
                box(61, y, 52, box_h, col["accent"], radius=6)
                c.setFont(font, 17)
                c.setFillColor(col["bg"] if theme.is_dark else col["white"])
                c.drawCentredString(87, y + box_h / 2 - 6, str(position + 1))

                box(122, y, 777, box_h, col["surface"], col["border"])
                c.setFont(font, 12.5)
                c.setFillColor(col["text"])
                c.drawString(142, y + box_h / 2 - 4.5, plain(step))
                y -= box_h + 9

            if spec.get("note"):
                c.setFont(font, 8.5)
                c.setFillColor(col["muted"])
                for line_index, line in enumerate(wrap(spec["note"], 8.5, 838)):
                    c.drawString(61, 62 - line_index * 11, line)

        elif kind == "screenshot":
            from reportlab.lib.utils import ImageReader

            image_path = PROJECT_ROOT / spec["image"]
            reader = ImageReader(str(image_path))
            iw, ih = reader.getSize()
            aspect = ih / iw

            max_w, max_h = 812, 322
            width = min(max_w, max_h / aspect)
            height = width * aspect
            x = 61 + (838 - width) / 2
            y = H - 122 - height

            c.setFillColor(col["surface"])
            c.setStrokeColor(col["border"])
            c.setLineWidth(0.8)
            c.roundRect(x - 5, y - 5, width + 10, height + 10, 6, stroke=1, fill=1)
            c.drawImage(reader, x, y, width=width, height=height)

            if spec.get("caption"):
                c.setFont(font, 9)
                c.setFillColor(col["muted"])
                for line_index, line in enumerate(wrap(spec["caption"], 9, 838)):
                    c.drawString(61, y - 22 - line_index * 12, line)

        elif kind == "compare":
            for side, (title_key, items_key, x0) in enumerate(
                [("left_title", "left", 61), ("right_title", "right", 511)]
            ):
                box(x0, H - 448, 388, 290, col["surface"], col["border"])
                c.setFillColor(col["wine"] if side == 0 else col["accent"])
                c.rect(x0, H - 161, 388, 3, stroke=0, fill=1)

                c.setFont(font, 12.5)
                c.setFillColor(
                    col["accent"] if (side or theme.is_dark) else col["wine"]
                )
                c.drawString(x0 + 20, H - 184, plain(spec[title_key]))

                y = H - 212
                for item in spec[items_key]:
                    for line_index, line in enumerate(wrap(f"·  {item}", 10, 348)):
                        c.setFont(font, 10)
                        c.setFillColor(col["text"])
                        c.drawString(x0 + 20 if line_index == 0 else x0 + 30, y, line)
                        y -= 13
                    y -= 6

            if spec.get("note"):
                c.setFont(font, 8.5)
                c.setFillColor(col["muted"])
                c.drawString(61, 58, plain(spec["note"]))

        footer(index)
        c.showPage()

    target.parent.mkdir(parents=True, exist_ok=True)
    c.save()


# ===========================================================================
# HTML
# ===========================================================================
def build_html(deck: dict, theme: Theme, target: Path) -> None:
    """Kendi kendine yeten, kaydırmalı HTML sunum üretir."""
    total = len(deck["slides"])
    labels = deck["labels"]
    body: list[str] = []

    for index, spec in enumerate(deck["slides"], start=1):
        kind = spec["type"]
        inner: list[str] = []

        if kind in {"title", "closing"}:
            items = ""
            if kind == "closing" and spec.get("items"):
                lis = "".join(f"<li>{to_html(i)}</li>" for i in spec["items"])
                items = f'<ul class="cover-list">{lis}</ul>'
            note = (
                f'<p class="cover-note">{to_html(spec["note"])}</p>'
                if spec.get("note") else ""
            )
            inner.append(
                f'<div class="cover">'
                f'<div class="cover-main">'
                f'<h1>{to_html(spec["title"])}</h1>'
                f'<p class="cover-sub">{to_html(spec["subtitle"])}</p>'
                f"{note}"
                f'<div class="rule"></div>'
                f'<p class="cover-meta">{to_html(spec.get("meta", ""))}</p>'
                f"</div>"
                f'<div class="cover-side">{items}</div>'
                f"</div>"
            )

        elif kind == "section":
            sub = (
                f'<p class="section-sub">{to_html(spec["subtitle"])}</p>'
                if spec.get("subtitle") else ""
            )
            inner.append(
                f'<div class="section">'
                f'<span class="section-count">{index} / {total}</span>'
                f'<div class="section-band"><h2>{to_html(spec["title"])}</h2>{sub}</div>'
                f"</div>"
            )

        else:
            sub = (
                f'<p class="s-sub">{to_html(spec["subtitle"])}</p>'
                if spec.get("subtitle") else ""
            )
            inner.append(
                f'<header class="s-head"><h2>{to_html(spec["title"])}</h2>{sub}'
                f'<div class="rule"></div></header>'
            )

            if kind == "bullets":
                lis = "".join(f"<li>{to_html(i)}</li>" for i in spec["items"])
                dense = " dense" if len(spec["items"]) > 6 else ""
                inner.append(f'<ul class="bullets{dense}">{lis}</ul>')

            elif kind == "stats":
                cards = "".join(
                    f'<div class="stat"><span class="stat-v">{html_mod.escape(v)}</span>'
                    f'<span class="stat-l">{html_mod.escape(label)}</span></div>'
                    for v, label in spec["items"]
                )
                cols = " four" if len(spec["items"]) > 6 else ""
                inner.append(f'<div class="stats{cols}">{cards}</div>')

            elif kind == "table":
                head = "".join(f"<th>{to_html(c)}</th>" for c in spec["columns"])
                rows = "".join(
                    "<tr>"
                    + "".join(
                        f'<td class="{"k" if i == 0 else ""}">{to_html(c)}</td>'
                        for i, c in enumerate(r)
                    )
                    + "</tr>"
                    for r in spec["rows"]
                )
                inner.append(
                    f'<div class="t-wrap"><table><thead><tr>{head}</tr></thead>'
                    f"<tbody>{rows}</tbody></table></div>"
                )

            elif kind == "flow":
                steps = "".join(
                    f'<li><span class="n">{i}</span><span class="s">{to_html(s)}</span></li>'
                    for i, s in enumerate(spec["steps"], start=1)
                )
                inner.append(f'<ol class="flow">{steps}</ol>')

            elif kind == "screenshot":
                import base64

                image_path = PROJECT_ROOT / spec["image"]
                encoded = base64.b64encode(image_path.read_bytes()).decode("ascii")
                caption = (
                    f'<p class="shot-cap">{to_html(spec["caption"])}</p>'
                    if spec.get("caption") else ""
                )
                inner.append(
                    f'<figure class="shot-wrap">'
                    f'<img class="shot" alt="{html_mod.escape(plain(spec["title"]))}" '
                    f'src="data:image/png;base64,{encoded}">'
                    f"{caption}</figure>"
                )

            elif kind == "compare":
                def panel(title, items, tone):
                    lis = "".join(f"<li>{to_html(i)}</li>" for i in items)
                    return (
                        f'<div class="panel {tone}"><h3>{to_html(title)}</h3>'
                        f"<ul>{lis}</ul></div>"
                    )

                inner.append(
                    '<div class="compare">'
                    + panel(spec["left_title"], spec["left"], "a")
                    + panel(spec["right_title"], spec["right"], "b")
                    + "</div>"
                )

            if spec.get("note"):
                inner.append(f'<p class="s-note">{to_html(spec["note"])}</p>')

        body.append(
            f'<section class="slide" data-n="{index}">'
            f'<div class="canvas">{"".join(inner)}'
            f'<footer class="s-foot"><span>{html_mod.escape(deck["footer"])}</span>'
            f"<span>{index}</span></footer></div></section>"
        )

    accent, bg, surface = f"#{theme.accent}", f"#{theme.bg}", f"#{theme.surface}"
    surface_alt, border = f"#{theme.surface_alt}", f"#{theme.border}"
    text, muted, wine = f"#{theme.text}", f"#{theme.muted}", f"#{theme.wine}"

    document = f"""<!doctype html>
<html lang="{deck['lang']}">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1, viewport-fit=cover">
<meta name="theme-color" content="{bg}">
<meta name="description" content="{html_mod.escape(deck['tagline'])}">
<title>{html_mod.escape(deck['deck_title'])}</title>
<style>
*{{box-sizing:border-box;margin:0;padding:0}}
:root{{
  --bg:{bg}; --surface:{surface}; --surface2:{surface_alt}; --border:{border};
  --text:{text}; --muted:{muted}; --accent:{accent}; --wine:{wine};
}}
html,body{{height:100%;background:var(--bg);color:var(--text);
  font-family:"DejaVu Sans","Source Sans 3",Inter,system-ui,sans-serif;
  -webkit-text-size-adjust:100%;overscroll-behavior:none}}
body{{display:flex;flex-direction:column;height:100dvh}}

header.top{{flex:0 0 auto;display:flex;align-items:center;gap:.55rem;
  padding:calc(.5rem + env(safe-area-inset-top)) .9rem .5rem;
  border-bottom:1px solid var(--border);background:var(--bg)}}
.brand{{font-size:.85rem;font-weight:600;white-space:nowrap;overflow:hidden;
  text-overflow:ellipsis}}
.ver{{flex:0 0 auto;font-size:.66rem;font-weight:700;color:{bg};
  background:var(--accent);border-radius:99px;padding:.14rem .5rem}}
.count{{margin-left:auto;font-size:.78rem;color:var(--muted);
  font-variant-numeric:tabular-nums}}
.bar{{flex:0 0 auto;height:2px;background:var(--border)}}
.bar>i{{display:block;height:100%;width:0;background:var(--accent);
  transition:width .18s ease}}

main{{flex:1 1 auto;display:flex;overflow-x:auto;overflow-y:hidden;
  scroll-snap-type:x mandatory;scroll-behavior:smooth;scrollbar-width:none;
  -webkit-overflow-scrolling:touch}}
main::-webkit-scrollbar{{display:none}}

.slide{{flex:0 0 100%;scroll-snap-align:center;scroll-snap-stop:always;
  display:flex;align-items:center;justify-content:center;padding:.6rem}}
.canvas{{width:100%;max-width:1180px;aspect-ratio:16/9;background:var(--bg);
  border:1px solid var(--border);border-radius:12px;position:relative;
  padding:clamp(1rem,3.1vw,2.5rem);display:flex;flex-direction:column;
  overflow:hidden}}

h1{{font-size:clamp(1.9rem,5.2vw,3.7rem);color:var(--accent);line-height:1.03;
  letter-spacing:-.02em}}
h2{{font-size:clamp(1.05rem,2.5vw,1.95rem);line-height:1.15;letter-spacing:-.01em}}
h3{{font-size:clamp(.82rem,1.5vw,1.05rem);color:var(--accent)}}
strong{{color:var(--accent);font-weight:600}}

.rule{{width:clamp(70px,9vw,124px);height:3px;background:var(--accent);
  margin:.7rem 0 0;border-radius:2px}}
.s-head{{flex:0 0 auto;margin-bottom:clamp(.7rem,2vw,1.5rem)}}
.s-sub{{color:var(--muted);font-size:clamp(.7rem,1.25vw,.95rem);margin-top:.35rem}}
.s-note{{margin-top:auto;padding-top:.7rem;color:var(--muted);
  font-size:clamp(.6rem,1.05vw,.78rem);line-height:1.45}}
.s-foot{{position:absolute;left:clamp(1rem,3.1vw,2.5rem);
  right:clamp(1rem,3.1vw,2.5rem);bottom:clamp(.55rem,1.6vw,1.05rem);
  display:flex;justify-content:space-between;color:var(--muted);
  font-size:clamp(.52rem,.85vw,.66rem)}}

/* Kapak */
.cover{{flex:1;display:grid;grid-template-columns:1.15fr .85fr;gap:1.5rem;
  align-items:center}}
.cover-sub{{font-size:clamp(.85rem,1.8vw,1.35rem);margin-top:.8rem;line-height:1.3}}
.cover-note{{color:var(--muted);font-size:clamp(.65rem,1.15vw,.88rem);
  margin-top:.55rem}}
.cover-meta{{color:var(--muted);font-size:clamp(.62rem,1.05vw,.8rem);
  margin-top:.7rem}}
.cover-side{{background:var(--surface);border-radius:12px;padding:1.2rem;
  align-self:stretch;display:flex;align-items:center}}
.cover-list{{list-style:none;display:grid;gap:.85rem}}
.cover-list li{{font-size:clamp(.66rem,1.15vw,.9rem);line-height:1.4;
  padding-left:1.1rem;position:relative}}
.cover-list li::before{{content:"◆";position:absolute;left:0;color:var(--accent)}}

/* Bölüm ayracı */
.section{{flex:1;display:flex;flex-direction:column;justify-content:center;
  position:relative}}
.section-count{{position:absolute;top:0;color:var(--muted);
  font-size:clamp(.6rem,1vw,.78rem)}}
.section-band{{background:var(--surface);border-left:4px solid var(--accent);
  border-radius:0 10px 10px 0;padding:clamp(1rem,2.6vw,2rem);
  margin:0 calc(-1 * clamp(1rem,3.1vw,2.5rem)) 0 0}}
.section-band h2{{font-size:clamp(1.4rem,3.4vw,2.6rem)}}
.section-sub{{color:var(--accent);font-size:clamp(.72rem,1.4vw,1.05rem);
  margin-top:.5rem}}

/* Maddeler */
.bullets{{list-style:none;display:flex;flex-direction:column;
  gap:clamp(.42rem,1.35vw,1rem);flex:1;justify-content:center}}
.bullets li{{position:relative;padding-left:1.15rem;line-height:1.42;
  font-size:clamp(.68rem,1.35vw,1.05rem)}}
.bullets.dense li{{font-size:clamp(.62rem,1.2vw,.95rem);line-height:1.36}}
.bullets li::before{{content:"";position:absolute;left:0;top:.52em;width:6px;
  height:6px;border-radius:2px;background:var(--accent)}}

/* Sayı kartları */
.stats{{display:grid;grid-template-columns:repeat(3,1fr);
  gap:clamp(.4rem,1.1vw,.85rem);flex:1;align-content:center}}
.stats.four{{grid-template-columns:repeat(4,1fr)}}
.stat{{background:var(--surface);border:1px solid var(--border);
  border-left:3px solid var(--accent);border-radius:9px;
  padding:clamp(.5rem,1.3vw,1rem);display:flex;flex-direction:column;gap:.2rem}}
.stat-v{{font-size:clamp(1.05rem,2.6vw,2rem);font-weight:700;color:var(--accent);
  line-height:1;font-variant-numeric:tabular-nums}}
.stat-l{{font-size:clamp(.54rem,.98vw,.76rem);color:var(--muted);line-height:1.3}}

/* Tablo */
.t-wrap{{flex:1;overflow:auto;display:flex;align-items:flex-start}}
table{{width:100%;border-collapse:separate;border-spacing:0 2px;
  font-size:clamp(.58rem,1.05vw,.85rem)}}
th{{background:var(--wine);color:#fff;text-align:left;font-weight:600;
  padding:clamp(.34rem,.85vw,.62rem) clamp(.42rem,1vw,.8rem)}}
th:first-child{{border-radius:6px 0 0 6px}}
th:last-child{{border-radius:0 6px 6px 0}}
td{{background:var(--surface);padding:clamp(.34rem,.85vw,.62rem)
  clamp(.42rem,1vw,.8rem);line-height:1.35;border-top:1px solid var(--border);
  border-bottom:1px solid var(--border)}}
tr:nth-child(even) td{{background:var(--surface2)}}
td:first-child{{border-left:1px solid var(--border);border-radius:6px 0 0 6px}}
td:last-child{{border-right:1px solid var(--border);border-radius:0 6px 6px 0}}
td.k{{color:var(--accent);font-weight:600}}

/* Akış */
.flow{{list-style:none;display:flex;flex-direction:column;
  gap:clamp(.32rem,.9vw,.62rem);flex:1;justify-content:center}}
.flow li{{display:flex;align-items:stretch;gap:.6rem}}
.flow .n{{flex:0 0 auto;width:clamp(1.75rem,3.4vw,2.9rem);
  display:grid;place-items:center;background:var(--accent);color:{bg};
  font-weight:700;border-radius:8px;font-size:clamp(.72rem,1.4vw,1.15rem)}}
.flow .s{{flex:1;background:var(--surface);border:1px solid var(--border);
  border-radius:8px;padding:clamp(.42rem,1.05vw,.82rem) clamp(.6rem,1.4vw,1.1rem);
  display:flex;align-items:center;line-height:1.32;
  font-size:clamp(.64rem,1.25vw,1rem)}}

/* Ekran görüntüsü */
.shot-wrap{{flex:1;display:flex;flex-direction:column;align-items:center;
  justify-content:center;gap:.55rem;min-height:0}}
.shot{{max-width:100%;max-height:100%;min-height:0;object-fit:contain;
  border:1px solid var(--border);border-radius:10px;
  box-shadow:0 4px 22px rgba(0,0,0,.28)}}
.shot-cap{{flex:0 0 auto;color:var(--muted);text-align:center;line-height:1.4;
  font-size:clamp(.6rem,1.1vw,.85rem);max-width:60rem}}

/* Karşılaştırma */
.compare{{display:grid;grid-template-columns:1fr 1fr;
  gap:clamp(.5rem,1.4vw,1.1rem);flex:1;align-content:stretch}}
.panel{{background:var(--surface);border:1px solid var(--border);border-radius:10px;
  padding:clamp(.65rem,1.5vw,1.2rem);display:flex;flex-direction:column;gap:.6rem;
  border-top:3px solid var(--wine)}}
.panel.b{{border-top-color:var(--accent)}}
.panel ul{{list-style:none;display:flex;flex-direction:column;
  gap:clamp(.3rem,.85vw,.6rem)}}
.panel li{{position:relative;padding-left:.85rem;line-height:1.38;
  font-size:clamp(.6rem,1.1vw,.88rem)}}
.panel li::before{{content:"·";position:absolute;left:.15rem;color:var(--accent);
  font-weight:700}}

nav.bottom{{flex:0 0 auto;display:flex;align-items:center;gap:.5rem;
  padding:.5rem .9rem calc(.5rem + env(safe-area-inset-bottom));
  border-top:1px solid var(--border);background:var(--bg)}}
nav button{{font:inherit;font-size:.78rem;color:var(--text);background:var(--surface);
  border:1px solid var(--border);border-radius:8px;padding:.42rem .85rem;
  cursor:pointer;min-height:34px}}
nav button:hover{{border-color:var(--accent)}}
nav button:focus-visible{{outline:2px solid var(--accent);outline-offset:2px}}
.hint{{margin:0 auto;font-size:.68rem;color:var(--muted)}}

@media (max-width:760px){{
  .cover{{grid-template-columns:1fr;gap:.8rem}}
  .cover-side{{display:none}}
  .stats,.stats.four{{grid-template-columns:repeat(2,1fr)}}
  .compare{{grid-template-columns:1fr;gap:.5rem}}
  .canvas{{aspect-ratio:auto;height:100%;overflow-y:auto}}
  .hint{{display:none}}
}}
@media print{{
  header.top,nav.bottom,.bar{{display:none}}
  body{{height:auto;display:block}}
  main{{display:block;overflow:visible}}
  .slide{{page-break-after:always;padding:0;display:block}}
  .canvas{{border:none;border-radius:0;max-width:none;height:100vh}}
}}
</style>
</head>
<body>

<header class="top">
  <span class="brand">{html_mod.escape(deck['product'])} — {html_mod.escape(deck['tagline'])}</span>
  <span class="ver">v{DECK_VERSION}</span>
  <span class="count"><b id="cur">1</b> {html_mod.escape(labels['of'])} {total}</span>
</header>
<div class="bar"><i id="prog"></i></div>

<main id="deck">
{"".join(body)}
</main>

<nav class="bottom">
  <button id="prev" type="button" aria-label="{html_mod.escape(labels['slide'])} −">←</button>
  <span class="hint">{html_mod.escape(labels['swipe'])}</span>
  <button id="next" type="button" aria-label="{html_mod.escape(labels['slide'])} +">→</button>
</nav>

<script>
(function () {{
  var deck = document.getElementById('deck');
  var slides = Array.prototype.slice.call(deck.querySelectorAll('.slide'));
  var cur = document.getElementById('cur');
  var prog = document.getElementById('prog');
  var total = slides.length;
  var index = 0;

  function paint() {{
    cur.textContent = index + 1;
    prog.style.width = ((index + 1) / total * 100) + '%';
    if (location.hash !== '#' + (index + 1)) {{
      history.replaceState(null, '', '#' + (index + 1));
    }}
  }}

  function go(n) {{
    index = Math.max(0, Math.min(total - 1, n));
    slides[index].scrollIntoView({{behavior: 'smooth', inline: 'center'}});
    paint();
  }}

  document.getElementById('prev').onclick = function () {{ go(index - 1); }};
  document.getElementById('next').onclick = function () {{ go(index + 1); }};

  document.addEventListener('keydown', function (e) {{
    if (e.key === 'ArrowRight' || e.key === 'PageDown' || e.key === ' ') {{
      e.preventDefault(); go(index + 1);
    }} else if (e.key === 'ArrowLeft' || e.key === 'PageUp') {{
      e.preventDefault(); go(index - 1);
    }} else if (e.key === 'Home') {{ go(0); }}
    else if (e.key === 'End') {{ go(total - 1); }}
  }});

  var ticking = false;
  deck.addEventListener('scroll', function () {{
    if (ticking) return;
    ticking = true;
    requestAnimationFrame(function () {{
      var n = Math.round(deck.scrollLeft / deck.clientWidth);
      if (n !== index) {{ index = Math.max(0, Math.min(total - 1, n)); paint(); }}
      ticking = false;
    }});
  }}, {{passive: true}});

  var start = parseInt((location.hash || '').slice(1), 10);
  if (start > 1 && start <= total) {{
    index = start - 1;
    slides[index].scrollIntoView({{behavior: 'auto', inline: 'center'}});
  }}
  paint();
}})();
</script>
</body>
</html>
"""
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_text(document, encoding="utf-8")


# ===========================================================================
# Ana akış
# ===========================================================================
DECK_VERSION = "0.1.0"

#: Çıktı klasörü. Sunum dosyaları depo kökünü kirletmez.
OUTPUT_DIR = PROJECT_ROOT / "docs" / "presentation"

#: ``_PUBLIC`` soneki, dosyanın herkese açık paylaşım için üretildiğini ve
#: yalnızca sentetik demo verisi içerdiğini belirtir.
OUTPUTS = [
    # (dil, tema, dosya kökü)
    ("tr", SCREEN, "Wine_House_Tanitim_PUBLIC"),
    ("tr", PRINT, "Wine_House_Tanitim_Baski_PUBLIC"),
    ("en", SCREEN, "Wine_House_Intro_EN_PUBLIC"),
    ("en", PRINT, "Wine_House_Intro_EN_Print_PUBLIC"),
]


def main() -> int:
    # Django ayarları (PDF yazı tipi yardımcısı için)
    import os

    sys.path.insert(0, str(PROJECT_ROOT / "src"))
    os.environ.setdefault("DJANGO_SETTINGS_MODULE", "winehouse.settings.dev")
    try:
        import django

        django.setup()
    except Exception:
        pass  # Yazı tipi bulucu Django olmadan da çalışır.

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    produced: list[tuple[str, int]] = []

    for lang, theme, stem in OUTPUTS:
        deck = DECKS[lang]

        pptx_path = OUTPUT_DIR / f"{stem}.pptx"
        build_pptx(deck, theme, pptx_path)
        produced.append((pptx_path.name, pptx_path.stat().st_size))

        pdf_path = OUTPUT_DIR / f"{stem}.pdf"
        build_pdf(deck, theme, pdf_path)
        produced.append((pdf_path.name, pdf_path.stat().st_size))

        # HTML yalnızca ekran teması için (baskı sürümü PDF olarak yeterli)
        if theme is SCREEN:
            html_path = OUTPUT_DIR / f"{stem}.html"
            build_html(deck, theme, html_path)
            produced.append((html_path.name, html_path.stat().st_size))

    print(f"{len(produced)} dosya üretildi:\n")
    for name, size in sorted(produced):
        print(f"  {name:<38} {size / 1024:8.0f} KB")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
